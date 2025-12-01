"""
Document service - file processing and document management.
"""
import os
import uuid
from pathlib import Path
from typing import List, Optional

import aiofiles
import structlog
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from models.document import Document

logger = structlog.get_logger()

# Supported file types
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".txt", ".md", ".pptx", ".ppt"}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB


class DocumentService:
    """Service for managing documents and file processing."""

    def __init__(self, upload_dir: str = None):
        """Initialize document service."""
        if upload_dir is None:
            # Detect AWS Elastic Beanstalk environment
            # Check for /var/app (staging or current) OR EB-specific env vars
            is_eb = (
                os.path.exists("/var/app") or
                os.environ.get("AWS_EXECUTION_ENV") or
                os.environ.get("ELASTIC_BEANSTALK_ENVIRONMENT_NAME")
            )
            if is_eb:
                # Running on AWS Elastic Beanstalk - use /tmp which is always writable
                upload_dir = "/tmp/uploads"
            else:
                # Local development or Docker
                upload_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads")
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    async def save_document(
        self,
        db: AsyncSession,
        category_id: int,
        filename: str,
        original_name: str,
        file_type: str,
        file_size: int,
        content: bytes,
        chapter: Optional[str] = None,
    ) -> Document:
        """
        Save an uploaded document.

        Args:
            db: Database session
            category_id: Category to associate with
            filename: Generated unique filename
            original_name: Original filename from upload
            file_type: File extension
            file_size: Size in bytes
            content: File content
            chapter: Optional chapter/topic name to associate with document

        Returns:
            Created Document model
        """
        # Generate storage path
        storage_path = str(self.upload_dir / filename)

        # Save file to disk
        async with aiofiles.open(storage_path, "wb") as f:
            await f.write(content)

        # Create database record
        document = Document(
            category_id=category_id,
            filename=filename,
            original_name=original_name,
            file_type=file_type,
            file_size=file_size,
            storage_path=storage_path,
            processed=False,
            chapter=chapter,
        )

        db.add(document)
        await db.flush()
        await db.refresh(document)

        logger.info(
            "document_saved",
            document_id=document.id,
            filename=original_name,
            file_type=file_type,
        )

        return document

    async def process_document(self, db: AsyncSession, document_id: int) -> Optional[str]:
        """
        Process a document and extract text content.

        Args:
            db: Database session
            document_id: Document ID

        Returns:
            Extracted text content
        """
        document = await self.get_document_by_id(db, document_id)
        if not document:
            return None

        try:
            content_text = await self._extract_text(document.storage_path, document.file_type)

            # Update document with extracted content
            document.content_text = content_text
            document.processed = True
            await db.flush()

            logger.info(
                "document_processed",
                document_id=document_id,
                text_length=len(content_text) if content_text else 0,
            )

            return content_text

        except Exception as e:
            logger.error("document_processing_error", document_id=document_id, error=str(e))
            raise

    async def _extract_text(self, file_path: str, file_type: str) -> str:
        """
        Extract text from a document based on file type.

        Args:
            file_path: Path to the file
            file_type: File extension

        Returns:
            Extracted text
        """
        file_type = file_type.lower()

        if file_type == ".pdf":
            return await self._extract_pdf_text(file_path)
        elif file_type in {".docx", ".doc"}:
            return await self._extract_word_text(file_path)
        elif file_type in {".pptx", ".ppt"}:
            return await self._extract_pptx_text(file_path)
        elif file_type in {".txt", ".md"}:
            return await self._extract_plain_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    async def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            # Use pypdf2 for PDF extraction
            from PyPDF2 import PdfReader

            reader = PdfReader(file_path)
            text_parts = []

            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)

            return "\n".join(text_parts)

        except Exception as e:
            logger.error("pdf_extraction_error", file_path=file_path, error=str(e))
            raise

    async def _extract_word_text(self, file_path: str) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)
            text_parts = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            return "\n".join(text_parts)

        except Exception as e:
            logger.error("word_extraction_error", file_path=file_path, error=str(e))
            raise

    async def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error("pptx_extraction_error", file_path=file_path, error=str(e))
            raise

    async def _extract_plain_text(self, file_path: str) -> str:
        """Extract text from plain text file."""
        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
            return await f.read()

    async def get_document_by_id(
        self,
        db: AsyncSession,
        document_id: int,
    ) -> Optional[Document]:
        """Get a document by ID."""
        result = await db.execute(
            select(Document).where(Document.id == document_id)
        )
        return result.scalar_one_or_none()

    async def get_documents_by_category(
        self,
        db: AsyncSession,
        category_id: int,
    ) -> List[Document]:
        """Get all documents for a category."""
        result = await db.execute(
            select(Document)
            .where(Document.category_id == category_id)
            .order_by(Document.created_at.desc())
        )
        return list(result.scalars().all())

    async def delete_document(
        self,
        db: AsyncSession,
        document_id: int,
    ) -> bool:
        """
        Delete a document and its file.

        Args:
            db: Database session
            document_id: Document ID

        Returns:
            True if deleted, False if not found
        """
        document = await self.get_document_by_id(db, document_id)
        if not document:
            return False

        # Delete file from disk
        try:
            if os.path.exists(document.storage_path):
                os.remove(document.storage_path)
        except Exception as e:
            logger.warning(
                "file_deletion_error",
                document_id=document_id,
                error=str(e),
            )

        # Delete database record
        await db.delete(document)
        await db.flush()

        logger.info("document_deleted", document_id=document_id)
        return True

    async def get_combined_content_for_category(
        self,
        db: AsyncSession,
        category_id: int,
        document_ids: Optional[List[int]] = None,
        chapter: Optional[str] = None,
    ) -> str:
        """
        Get combined text content from documents in a category.

        Args:
            db: Database session
            category_id: Category ID
            document_ids: Optional list of specific document IDs
            chapter: Optional chapter to filter by - only includes documents
                     that match this chapter

        Returns:
            Combined text content
        """
        query = select(Document).where(Document.category_id == category_id).where(Document.processed == True)

        if document_ids:
            query = query.where(Document.id.in_(document_ids))

        if chapter:
            query = query.where(Document.chapter == chapter)

        result = await db.execute(query)
        documents = result.scalars().all()
        content_parts = [doc.content_text for doc in documents if doc.content_text]

        return "\n\n---\n\n".join(content_parts)

    async def get_chapters_for_category(
        self,
        db: AsyncSession,
        category_id: int,
    ) -> List[str]:
        """
        Get all unique chapters for documents in a category.

        Args:
            db: Database session
            category_id: Category ID

        Returns:
            List of unique chapter names
        """
        result = await db.execute(
            select(Document.chapter)
            .where(Document.category_id == category_id)
            .where(Document.chapter.isnot(None))
            .distinct()
        )
        chapters = [row[0] for row in result.all() if row[0]]
        return sorted(chapters)

    async def update_document_chapter(
        self,
        db: AsyncSession,
        document_id: int,
        chapter: Optional[str],
    ) -> Optional[Document]:
        """
        Update the chapter field for a document.

        Args:
            db: Database session
            document_id: Document ID
            chapter: New chapter name (or None to clear)

        Returns:
            Updated document or None if not found
        """
        document = await self.get_document_by_id(db, document_id)
        if not document:
            return None

        document.chapter = chapter
        await db.flush()
        await db.refresh(document)

        logger.info("document_chapter_updated", document_id=document_id, chapter=chapter)
        return document

    def generate_filename(self, original_name: str) -> str:
        """Generate a unique filename for storage."""
        ext = Path(original_name).suffix.lower()
        unique_id = str(uuid.uuid4())[:8]
        return f"{unique_id}_{original_name}"

    def validate_file(self, filename: str, file_size: int) -> tuple[bool, str]:
        """
        Validate uploaded file.

        Returns:
            Tuple of (is_valid, error_message)
        """
        ext = Path(filename).suffix.lower()

        if ext not in SUPPORTED_EXTENSIONS:
            return False, f"Unsupported file type: {ext}. Supported: {', '.join(SUPPORTED_EXTENSIONS)}"

        if file_size > MAX_FILE_SIZE:
            return False, f"File too large. Maximum size: {MAX_FILE_SIZE // (1024 * 1024)}MB"

        return True, ""


# Global service instance
document_service = DocumentService()
